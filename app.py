
import os
from flask import Flask, request, render_template, redirect, url_for
from pptx import Presentation
from openai import OpenAI
from markdown import markdown
from PyPDF2 import PdfReader
import fitz  # PyMuPDF

client = OpenAI(api_key="openai_key")

app = Flask(__name__)
app.secret_key = "fitz_secret"
UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# In-memory storage
app.config['DATA'] = {}

def extract_pdf_text(pdf_path):
    reader = PdfReader(pdf_path)
    slides = []
    for page in reader.pages:
        text = page.extract_text()
        slides.append(text.strip() if text else "")
    return slides

def render_pdf_to_images(pdf_path, output_folder):
    os.makedirs(output_folder, exist_ok=True)
    doc = fitz.open(pdf_path)
    image_paths = []
    for i, page in enumerate(doc):
        pix = page.get_pixmap(dpi=150)
        img_path = os.path.join(output_folder, f"slide_{i}.png")
        pix.save(img_path)
        image_paths.append(f"/static/slides/slide_{i}.png")
    return image_paths

def extract_slide_text(ppt_path):
    prs = Presentation(ppt_path)
    slides = []
    for slide in prs.slides:
        content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content.append(shape.text)
        slides.append('\n'.join(content))
    return slides

def generate_explanation(text, mode):
    prompt = (
        f"Summarize this slide in 3-4 sentences:\n\n{text}"
        if mode == "summary"
        else f"Explain the following slide in detail and give a real-world example:\n\n{text}"
    )
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7
    )
    return markdown(response.choices[0].message.content)

def generate_flashcards(text):
    prompt = (
        "Generate flashcards (Q&A pairs) based on the following content. "
        "Create as many as are useful to understand the material clearly. "
        "Format each flashcard clearly with a question and answer:\n\n"
        f"{text}"
    )
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7
    )
    return markdown(response.choices[0].message.content)

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        app.config['DATA'].clear()

        file = request.files['ppt']
        mode = request.form.get("mode")
        filename = file.filename
        path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(path)

        if filename.endswith(".pptx"):
            slides = extract_slide_text(path)
            images = []
        elif filename.endswith(".pdf"):
            slides = extract_pdf_text(path)
            images = render_pdf_to_images(path, "static/slides")
        else:
            return "Only .pptx and .pdf files are supported."

        app.config['DATA']['slides'] = slides
        app.config['DATA']['images'] = images
        app.config['DATA']['questions'] = [''] * len(slides)
        app.config['DATA']['answers'] = [''] * len(slides)
        app.config['DATA']['explanations'] = [''] * len(slides)
        app.config['DATA']['flashcards'] = [''] * len(slides)
        app.config['DATA']['mode'] = mode

        return redirect(url_for("slide", slide_id=0))
    return render_template("index.html")

@app.route("/slide/<int:slide_id>")
def slide(slide_id):
    slides = app.config['DATA'].get('slides', [])
    if slide_id < 0 or slide_id >= len(slides):
        return redirect(url_for("index"))

    explanations = app.config['DATA']['explanations']
    flashcards = app.config['DATA']['flashcards']

    if not explanations[slide_id]:
        text = slides[slide_id]
        mode = app.config['DATA'].get('mode', 'summary')
        explanations[slide_id] = generate_explanation(text, mode)
        flashcards[slide_id] = generate_flashcards(text)

    return render_template(
        "slide.html",
        slide_id=slide_id,
        image=app.config['DATA']['images'][slide_id] if app.config['DATA'].get('images') else None,
        explanation=explanations[slide_id],
        flashcard=flashcards[slide_id],
        question=app.config['DATA']['questions'][slide_id],
        answer=app.config['DATA']['answers'][slide_id],
        total=len(slides)
    )

@app.route("/ask", methods=["POST"])
def ask():
    slide_id = int(request.form['slide_index'])
    question = request.form['question']
    context = app.config['DATA']['slides'][slide_id]

    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are an AI presenter answering questions about a slide."},
            {"role": "user", "content": f"Slide content:\n{context}\n\nQuestion:\n{question}"}
        ]
    )
    answer = markdown(response.choices[0].message.content)

    app.config['DATA']['questions'][slide_id] = question
    app.config['DATA']['answers'][slide_id] = answer

    return redirect(url_for("slide", slide_id=slide_id))

if __name__ == "__main__":
    app.run(debug=True)

